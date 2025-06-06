import streamlit as st
import os
from openai import OpenAI
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# Inicializar cliente OpenAI
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

# Configuración visual de la app 
st.set_page_config(page_title="ChatBot de clases", page_icon="🔖", layout="centered")

st.markdown("""
    <style>
    body {
        background-color: #fff0f5;
    }
    .main {
        background-color: #ffe6f0;
        padding: 2rem;
        border-radius: 12px;
    }
    h1, h2, h3, h4 {
        color: #c2185b;
    }
    .stButton button {
        background-color: #ec407a;
        color: white;
    }
    </style>
""", unsafe_allow_html=True)

# Logo institucional
st.image("logo uach.png", width=100)

# Banner institucional
st.markdown("""
<div style="background-color:#ffc0cb;padding:20px;border-radius:10px;text-align:center">
    <h1 style="color:#880e4f;">ChatBot de clases</h1>
    <h3 style="color:#6a1b9a;">Facultad de Medicina y Ciencias Biomédicas</h3>
    <h4 style="color:#6a1b9a;">Universidad Autónoma de Chihuahua</h4>
</div>
""", unsafe_allow_html=True)

# Rutas de archivos fuente
pptx_path = "clase_001_algo.pptx"
txt_path = "capitulo_hemodinamia_guyton.txt"
video_url_base = "https://youtu.be/kdjhkfehkurfhckehALGO"

# Índice temático del video
temas_video = {
    "TEMA A": "7:45",
    "TEMA B": "8:59",
    "TEMA C": "12:11",
    "TEMA D": "12:57",
    "TEMA E": "14:22"
}

# Funciones para extracción de texto
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    return [" ".join(shape.text for shape in slide.shapes if hasattr(shape, "text")).strip() for slide in prs.slides]

def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return [p.strip() for p in f.read().split("\n\n") if len(p.strip()) > 60]

# Cargar documentos en memoria
slides = extract_text_from_pptx(pptx_path)
chapter = extract_text_from_txt(txt_path)
all_docs = slides + chapter

# Entrada del usuario
query = st.text_input("🎓 Escribe tu pregunta:")

if query:
    # Buscar si hay tema del video relacionado
    query_lower = query.lower()
    tema_encontrado = next(((t, m) for t, m in temas_video.items() if t in query_lower), None)

    # Obtener contexto textual relevante
    vectorizer = TfidfVectorizer().fit_transform([query] + all_docs)
    similarity = cosine_similarity(vectorizer[0:1], vectorizer[1:])
    top_indices = similarity[0].argsort()[-3:][::-1]
    context = "\n\n".join([all_docs[i] for i in top_indices])

    # Crear prompt para GPT
    prompt = f"""
Eres un asistente de clases de la Facultad de Medicina de la UACH.
Responde únicamente con base en los siguientes materiales audiovisuales
No inventes ni agregues información externa.

PREGUNTA: {query}

MATERIALES:
{context}

RESPUESTA:
"""

    # Llamada a GPT-4
    with st.spinner("🤖 Consultando a GPT-4..."):
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2,
            max_tokens=600
        )

    # Mostrar respuesta en pantalla
    st.subheader(":open_book: Respuesta elaborada con base en tus clases:")
    st.write(response.choices[0].message.content)
    st.caption("📚 Esta respuesta se generó con base en tus presentaciones, lectura y video. No se utilizó información externa.")

    # Mostrar enlace al minuto del video si aplica
    if tema_encontrado:
        st.markdown(f"🎯 **Puedes encontrar la explicación de *{tema_encontrado[0]}* en el minuto {tema_encontrado[1]} del video.**")
        minutos, segundos = map(int, tema_encontrado[1].split(":"))
        start_time = minutos * 60 + segundos
        video_url = f"{video_url_base}&start={start_time}"
    else:
        video_url = video_url_base

    # Reproductor de video embebido
    st.markdown("**🎥 Explicación en video:**")
    st.video(video_url)
